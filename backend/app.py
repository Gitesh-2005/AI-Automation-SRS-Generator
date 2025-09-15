import os
import asyncio
from pathlib import Path
from typing import List, Optional, Dict, Any, AsyncGenerator, Tuple
from fastapi import FastAPI, UploadFile, File, HTTPException, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel
import aiofiles
import uuid
from datetime import datetime, timedelta
import json
import socketio 
import easyocr 
from langchain_core.messages import HumanMessage, AIMessage, SystemMessage, BaseMessage
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.output_parsers import StrOutputParser
from langgraph.checkpoint.memory import MemorySaver
from langgraph.graph import StateGraph, MessagesState, END
from langchain_core.load.dump import dumps as langchain_dumps
from langchain_core.load.load import loads as langchain_loads
from typing_extensions import Annotated
import PyPDF2
import docx
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from dotenv import load_dotenv
import logging
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type
from groq import RateLimitError
import hashlib
from collections import deque
from time import time
import copy
import re
from asyncio import TimeoutError
from contextlib import asynccontextmanager
import markdown
from markdown.extensions import Extension
from markdown.treeprocessors import Treeprocessor

try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Directory setup
UPLOAD_DIR = Path("uploads")
OUTPUT_DIR = Path("output")
TEMP_DIR = Path("temp")

for directory in [UPLOAD_DIR, OUTPUT_DIR, TEMP_DIR]:
    directory.mkdir(exist_ok=True)

# Initialize OCR
try:
    ocr = easyocr.Reader(['en'])
    logger.info("EasyOCR initialized successfully")
except Exception as e:
    logger.warning(f"Failed to initialize EasyOCR: {e}")
    ocr = None

load_dotenv()

API_KEYS = os.getenv("API_KEY")

llms: List[Tuple[ChatGroq, str]] = []
api_calls_per_key: Dict[str, deque] = {}
current_llm_index = 0
MAX_REQUESTS_PER_MINUTE = 25
llm_cache: Dict[str, str] = {}
TOKEN_LIMIT = 6000  # Token limit for llama-3.1-8b-instant
api_key_cooldowns: Dict[str, float] = {}  # Tracks cooldown expiration time for each API key

def estimate_tokens(text: str) -> int:
    """Estimate token count (approximate: 1 token ~ 4 characters)."""
    return len(text) // 4 + 1

def chunk_input(text: str, max_tokens: int = 5000) -> List[str]:
    """Split input text into token-safe chunks."""
    words = text.split()
    chunks = []
    current_chunk = []
    current_tokens = 0

    for word in words:
        word_tokens = estimate_tokens(word)
        if current_tokens + word_tokens > max_tokens:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_tokens = word_tokens
        else:
            current_chunk.append(word)
            current_tokens += word_tokens

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks

def format_streaming_content(content: str) -> str:
    """Enhanced formatting for streaming content with proper markdown rendering."""
    if not content:
        return content
    
    # Clean and format the content
    formatted = content
    
    # Ensure proper spacing after sentences
    formatted = re.sub(r'([.!?])([A-Z])', r'\1 \2', formatted)
    
    # Fix heading formatting
    formatted = re.sub(r'^(#{1,6})([^\s])', r'\1 \2', formatted, flags=re.MULTILINE)
    
    # Ensure proper list formatting
    formatted = re.sub(r'^([0-9]+\.)([^\s])', r'\1 \2', formatted, flags=re.MULTILINE)
    formatted = re.sub(r'^([-*+])([^\s])', r'\1 \2', formatted, flags=re.MULTILINE)
    
    # Add line breaks after section headers
    formatted = re.sub(r'(^#{1,6}.*$)', r'\1\n', formatted, flags=re.MULTILINE)
    
    # Ensure double line breaks between major sections
    formatted = re.sub(r'\n(#{1,2}\s)', r'\n\n\1', formatted)
    
    # Clean up excessive whitespace but preserve intentional formatting
    formatted = re.sub(r'\n{3,}', '\n\n', formatted)
    
    return formatted

def initialize_groq_llms():
    global llms, api_calls_per_key, api_key_cooldowns
    llms = []
    api_calls_per_key = {}
    api_key_cooldowns = {}
    for idx, api_key in enumerate(API_KEYS):
        if not api_key:
            logger.warning(f"API key {idx+1} not provided")
            continue
        try:
            llm = ChatGroq(
                model="llama-3.1-8b-instant",
                api_key=api_key,
                temperature=0.1,
                max_tokens=4096
            )
            llms.append((llm, api_key))
            api_calls_per_key[api_key] = deque(maxlen=60)
            api_key_cooldowns[api_key] = 0
            logger.info(f"Groq LLM initialized for API key {idx+1}")
        except Exception as e:
            logger.error(f"Failed to initialize Groq LLM for API key {idx+1}: {e}")
    return len(llms) > 0

def get_next_llm() -> Optional[Tuple[ChatGroq, str]]:
    global current_llm_index
    if not llms:
        return None
    
    original_index = current_llm_index
    current_time = time()
    attempts = 0
    max_attempts = len(llms) * 2  # Allow extra loops for cooldowns
    
    while attempts < max_attempts:
        llm_tuple = llms[current_llm_index]
        api_key = llm_tuple[1]
        
        if api_key_cooldowns.get(api_key, 0) > current_time:
            logger.info(f"API key {api_key[-4:]} is in cooldown until {api_key_cooldowns[api_key]}")
            current_llm_index = (current_llm_index + 1) % len(llms)
            attempts += 1
            continue
        
        api_calls = api_calls_per_key.get(api_key, deque(maxlen=60))
        while api_calls and current_time - api_calls[0] > 60:
            api_calls.popleft()
        
        if len(api_calls) < MAX_REQUESTS_PER_MINUTE:
            api_calls.append(current_time)
            next_index = (current_llm_index + 1) % len(llms)
            current_llm_index = next_index
            return llm_tuple
        
        logger.warning(f"Rate limit hit for API key {api_key[-4:]}")
        current_llm_index = (current_llm_index + 1) % len(llms)
        attempts += 1
    
    # If all failed, wait for the earliest cooldown to expire
    min_cooldown = min((api_key_cooldowns.get(llms[i][1], 0) for i in range(len(llms))), default=0)
    if min_cooldown > current_time:
        sleep_time = min_cooldown - current_time + 1  # +1 second buffer
        logger.info(f"All API keys unavailable, waiting {sleep_time:.2f} seconds")
        asyncio.create_task(asyncio.sleep(sleep_time))  # Non-blocking wait in loop context
    
    # Fallback to first available after wait
    llm_tuple = llms[0]
    api_calls_per_key[llm_tuple[1]].append(current_time)
    current_llm_index = 1
    return llm_tuple

def get_cache_key(input_dict: dict) -> str:
    def serialize_for_cache(obj):
        if isinstance(obj, BaseMessage):
            return {
                "type": obj.__class__.__name__,
                "content": str(obj.content) if obj.content else "",
                "additional_kwargs": getattr(obj, "additional_kwargs", {}),
                "id": getattr(obj, "id", None)
            }
        elif isinstance(obj, list) and obj and isinstance(obj[0], BaseMessage):
            return [serialize_for_cache(item) for item in obj]
        elif isinstance(obj, datetime):
            return obj.isoformat()
        elif isinstance(obj, dict):
            return {k: serialize_for_cache(v) for k, v in obj.items()}
        elif isinstance(obj, list):
            return [serialize_for_cache(item) for item in obj]
        else:
            return obj

    serializable_dict = {k: serialize_for_cache(v) for k, v in input_dict.items()}
    input_str = json.dumps(serializable_dict, sort_keys=True)
    return hashlib.md5(input_str.encode('utf-8')).hexdigest()

@retry(
    stop=stop_after_attempt(15),  # Increased attempts for completion guarantee
    wait=wait_exponential(multiplier=1, min=2, max=60),
    retry=retry_if_exception_type((RateLimitError, Exception))  # Retry on any exception
)
async def invoke_llm_with_retry(llm: ChatGroq, prompt: ChatPromptTemplate, output_parser, input_dict: dict) -> str:
    cache_key = get_cache_key(input_dict)
    if cache_key in llm_cache:
        logger.info(f"Cache hit for key: {cache_key}")
        return llm_cache[cache_key]

    # Track attempts across all API keys to ensure completion
    total_attempts = 0
    max_total_attempts = len(API_KEYS) * 5  # Allow multiple cycles through all keys
    
    while total_attempts < max_total_attempts:
        total_attempts += 1
        llm_tuple = get_next_llm()
        
        if not llm_tuple:
            # If no LLM available, wait and try again
            logger.warning(f"No LLM available on attempt {total_attempts}, waiting...")
            await asyncio.sleep(5)
            continue
            
        current_llm, api_key = llm_tuple
        current_time = time()
        api_calls = api_calls_per_key.get(api_key, deque(maxlen=60))

        # Clean old API calls
        while api_calls and current_time - api_calls[0] > 60:
            api_calls.popleft()

        # Check rate limit
        if len(api_calls) >= MAX_REQUESTS_PER_MINUTE:
            logger.warning(f"Rate limit approaching for API key {api_key[-4:]} on attempt {total_attempts}")
            await asyncio.sleep(2)
            continue

        api_calls.append(current_time)
        logger.info(f"Making API call (attempt {total_attempts}/{max_total_attempts}) with API key {api_key[-4:]}")

        try:
            chain = prompt | current_llm | output_parser
            response = await chain.ainvoke(input_dict)
            
            # Validate response before caching
            if response and len(str(response).strip()) > 10:
                llm_cache[cache_key] = response
                logger.info(f"API call completed successfully on attempt {total_attempts}")
                return response
            else:
                logger.warning(f"Received invalid/empty response on attempt {total_attempts}, retrying...")
                continue
                
        except RateLimitError as e:
            headers = getattr(e, 'response', {}).headers if hasattr(e, 'response') else {}
            retry_after = float(headers.get("retry-after", 60))
            api_key_cooldowns[api_key] = current_time + retry_after
            logger.warning(f"Rate limit hit for API key {api_key[-4:]} on attempt {total_attempts}; cooldown until {api_key_cooldowns[api_key]}")
            await asyncio.sleep(min(retry_after, 30))  # Cap wait time at 30 seconds
            continue
            
        except Exception as e:
            logger.error(f"Unexpected error in LLM invocation with API key {api_key[-4:]} on attempt {total_attempts}: {e}")
            await asyncio.sleep(2)
            continue
    
    # If we've exhausted all attempts, return a fallback response
    logger.error(f"Failed to get response after {total_attempts} attempts across all API keys")
    return "Failed to generate response. Please try again later."

@retry(
    stop=stop_after_attempt(10),
    wait=wait_exponential(multiplier=1, min=2, max=60),
    retry=retry_if_exception_type((RateLimitError, Exception))
)
async def invoke_llm_stream_with_retry(llm: ChatGroq, prompt: ChatPromptTemplate, output_parser, input_dict: dict) -> AsyncGenerator[str, None]:
    cache_key = get_cache_key(input_dict)
    if cache_key in llm_cache:
        logger.info(f"Cache hit for key: {cache_key}")
        cached_response = llm_cache[cache_key]
        # Format cached response with proper markdown breaks
        formatted = re.sub(r'\n\n+', '\n\n', cached_response)
        words = formatted.split()
        for word in words:
            yield word + (" " if not word.endswith(('\n', '\n\n')) else "")
            await asyncio.sleep(0.05)
        return

    while True:  # Loop until success
        llm_tuple = get_next_llm()
        if not llm_tuple:
            yield "No LLM available. Retrying..."
            await asyncio.sleep(5)
            continue
        current_llm, api_key = llm_tuple

        current_time = time()
        api_calls = api_calls_per_key.get(api_key, deque(maxlen=60))

        while api_calls and current_time - api_calls[0] > 60:
            api_calls.popleft()

        while len(api_calls) >= MAX_REQUESTS_PER_MINUTE:
            logger.warning(f"Rate limit approaching for API key {api_key[-4:]}")
            await asyncio.sleep(5)
            current_time = time()
            while api_calls and current_time - api_calls[0] > 60:
                api_calls.popleft()

        api_calls.append(current_time)
        logger.info(f"Making streaming API call with input keys: {input_dict.keys()} using API key {api_key[-4:]}")

        full_input = "\n".join([f"{k}: {v}" for k, v in input_dict.items() if isinstance(v, str)])
        if estimate_tokens(full_input) > TOKEN_LIMIT:
            logger.warning(f"Input exceeds token limit ({estimate_tokens(full_input)} > {TOKEN_LIMIT}). Streaming chunks...")
            chunks = chunk_input(full_input, max_tokens=TOKEN_LIMIT - 1000)
            full_response = ""
            for i, chunk in enumerate(chunks):
                logger.info(f"Streaming chunk {i+1}/{len(chunks)}")
                chunk_input = {k: chunk if k == list(input_dict.keys())[0] else v for k, v in input_dict.items()}
                chain = prompt | current_llm | output_parser
                try:
                    async for partial_chunk in chain.astream(chunk_input):
                        if partial_chunk:
                            # Enhanced formatting with proper markdown rendering
                            formatted_chunk = format_streaming_content(partial_chunk)
                            full_response += partial_chunk
                            yield formatted_chunk
                            await asyncio.sleep(0.1)  # Faster streaming for better UX
                    yield "\n\n✓ **Chunk Complete**\n\n"
                    await asyncio.sleep(2)
                except RateLimitError as e:
                    headers = getattr(e, 'response', {}).headers if hasattr(e, 'response') else {}
                    retry_after = float(headers.get("retry-after", 60))
                    api_key_cooldowns[api_key] = current_time + retry_after
                    logger.warning(f"Rate limit hit for API key {api_key[-4:]}; cooldown until {api_key_cooldowns[api_key]}")
                    yield f"Switching API key due to rate limit...\n\n"
                    # Switch and retry chunk
                    continue
                except TimeoutError:
                    logger.error(f"Timeout while streaming chunk {i+1} with API key {api_key[-4:]}")
                    yield f"Error: Stream timeout for chunk {i+1}\n\n"
                    return
                except Exception as e:
                    logger.error(f"Unexpected error in streaming chunk {i+1} with API key {api_key[-4:]}: {e}")
                    yield f"Error: {str(e)}\n\n"
                    return
            llm_cache[cache_key] = full_response
            logger.info("Streaming API call completed successfully")
            return

        try:
            chain = prompt | current_llm | output_parser
            full_response = ""
            async for chunk in chain.astream(input_dict):
                if chunk:
                    # Enhanced formatting for better readability
                    formatted_chunk = format_streaming_content(chunk)
                    full_response += chunk
                    yield formatted_chunk
                    await asyncio.sleep(0.1)  # Faster streaming
            llm_cache[cache_key] = full_response
            logger.info("Streaming API call completed successfully")
            return
        except RateLimitError as e:
            headers = getattr(e, 'response', {}).headers if hasattr(e, 'response') else {}
            retry_after = float(headers.get("retry-after", 60))
            api_key_cooldowns[api_key] = current_time + retry_after
            logger.warning(f"Rate limit hit for API key {api_key[-4:]}; cooldown until {api_key_cooldowns[api_key]}")
            yield f"Switching API key due to rate limit...\n\n"
            # Continue loop to switch key
        except TimeoutError:
            logger.error(f"Timeout in LLM streaming invocation with API key {api_key[-4:]}")
            yield "Error: Stream timeout\n\n"
            await asyncio.sleep(5)
        except Exception as e:
            logger.error(f"Unexpected error in LLM streaming invocation with API key {api_key[-4:]}: {e}")
            yield f"Error: {str(e)}\n\n"
            await asyncio.sleep(2)

def safe_json_serialize(data: Any) -> Any:
    if isinstance(data, BaseMessage):
        return {
            "type": data.__class__.__name__,
            "content": str(data.content) if data.content else "",
            "additional_kwargs": getattr(data, "additional_kwargs", {}),
            "id": getattr(data, "id", None)
        }
    elif isinstance(data, list) and data and isinstance(data[0], BaseMessage):
        return [safe_json_serialize(msg) for msg in data]
    elif isinstance(data, datetime):
        return data.isoformat()
    elif isinstance(data, dict):
        return {k: safe_json_serialize(v) for k, v in data.items()}
    elif isinstance(data, list):
        return [safe_json_serialize(item) for item in data]
    else:
        return data

class AgenticSRSState(MessagesState):
    messages: Annotated[list, MessagesPlaceholder("messages")]
    requirements_analysis: Optional[str] = None
    functional_requirements: Optional[str] = None
    non_functional_requirements: Optional[str] = None
    system_architecture: Optional[str] = None
    validation_feedback: Optional[str] = None
    final_srs: Optional[str] = None
    user_input: str = ""
    edit_context: Optional[str] = None
    edit_instructions: Optional[str] = None
    current_content: Optional[str] = None

async def analyze_requirements_node(state: AgenticSRSState) -> Dict[str, Any]:
    llm_tuple = get_next_llm()
    if not llm_tuple:
        return {"requirements_analysis": "LLM not available"}
    llm, _ = llm_tuple

    document_context = state.get("document_context", "")
    user_input = state["user_input"]
    
    # Format document context as reference material, not as primary content
    context_prompt = ""
    if document_context:
        if isinstance(document_context, list) and document_context:
            context_summaries = []
            for ctx in document_context[:3]:  # Limit to top 3 most relevant
                if isinstance(ctx, dict):
                    filename = ctx.get('filename', 'Document')
                    analysis = ctx.get('analysis', 'No analysis available')
                    context_summaries.append(f"Reference - {filename}: {analysis[:300]}...")
            context_prompt = "\n\nReference Documents (for context only):\n" + "\n".join(context_summaries)
        elif isinstance(document_context, str) and document_context.strip():
            context_prompt = f"\n\nReference Context: {document_context[:500]}..."
    
    prompt = ChatPromptTemplate.from_messages([
        ("system", """You are analyzing USER REQUIREMENTS to create an SRS. The user input is PRIMARY and contains the actual requirements to analyze. Any reference documents provided are ONLY for context and domain understanding - do not generate SRS based on the reference documents.
        
        CRITICAL INSTRUCTIONS:
        1. Focus PRIMARILY on the user's input requirements
        2. Use reference documents ONLY to understand domain context, terminology, and best practices
        3. Do NOT generate requirements directly from reference documents
        4. Extract business objectives, functional needs, and constraints from USER INPUT
        5. Use document context to inform your analysis style and terminology
        
        Output analysis in markdown with sections: Business Objectives, Key Functions, Constraints, Domain Context."""),
        MessagesPlaceholder("messages"),
        ("human", "USER REQUIREMENTS TO ANALYZE:\n{user_input}{context_prompt}\n\nAnalyze the USER REQUIREMENTS above (use reference context for understanding only):")
    ])
    parser = StrOutputParser()

    analysis = await invoke_llm_with_retry(llm, prompt, parser, {
        "user_input": user_input,
        "context_prompt": context_prompt,
        "messages": state["messages"]
    })
    return {
        "requirements_analysis": analysis,
        "messages": state["messages"] + [HumanMessage(content=analysis)]
    }

async def generate_functional_requirements_node(state: AgenticSRSState) -> Dict[str, Any]:
    llm_tuple = get_next_llm()
    if not llm_tuple:
        return {"functional_requirements": "LLM not available"}
    llm, _ = llm_tuple

    prompt = ChatPromptTemplate.from_messages([
        ("system", """Generate detailed functional requirements based on the analysis. Use numbered lists and use cases. Output in markdown under 'Functional Requirements'."""),
        MessagesPlaceholder("messages"),
        ("human", "{analysis}")
    ])
    parser = StrOutputParser()

    functional = await invoke_llm_with_retry(llm, prompt, parser, {
        "analysis": state["requirements_analysis"],
        "messages": state["messages"]
    })
    return {
        "functional_requirements": functional,
        "messages": state["messages"] + [HumanMessage(content=functional)]
    }

async def generate_non_functional_requirements_node(state: AgenticSRSState) -> Dict[str, Any]:
    llm_tuple = get_next_llm()
    if not llm_tuple:
        return {"non_functional_requirements": "LLM not available"}
    llm, _ = llm_tuple

    prompt = ChatPromptTemplate.from_messages([
        ("system", """Derive non-functional requirements like performance, security, usability from the analysis and functional requirements. Output in markdown under 'Non-Functional Requirements'."""),
        MessagesPlaceholder("messages"),
        ("human", "{analysis}\n\n{functional}")
    ])
    parser = StrOutputParser()

    non_functional = await invoke_llm_with_retry(llm, prompt, parser, {
        "analysis": state["requirements_analysis"],
        "functional": state["functional_requirements"],
        "messages": state["messages"]
    })
    return {
        "non_functional_requirements": non_functional,
        "messages": state["messages"] + [HumanMessage(content=non_functional)]
    }

async def generate_system_architecture_node(state: AgenticSRSState) -> Dict[str, Any]:
    llm_tuple = get_next_llm()
    if not llm_tuple:
        return {"system_architecture": "LLM not available"}
    llm, _ = llm_tuple

    prompt = ChatPromptTemplate.from_messages([
        ("system", """Outline a high-level system architecture including components, data flow, and tech stack suggestions based on all prior sections. Output in markdown under 'System Architecture'."""),
        MessagesPlaceholder("messages"),
        ("human", "{analysis}\n\n{functional}\n\n{non_functional}")
    ])
    parser = StrOutputParser()

    architecture = await invoke_llm_with_retry(llm, prompt, parser, {
        "analysis": state["requirements_analysis"],
        "functional": state["functional_requirements"],
        "non_functional": state["non_functional_requirements"],
        "messages": state["messages"]
    })
    return {
        "system_architecture": architecture,
        "messages": state["messages"] + [HumanMessage(content=architecture)]
    }

async def validate_and_refine_node(state: AgenticSRSState) -> Dict[str, Any]:
    llm_tuple = get_next_llm()
    if not llm_tuple:
        return {"validation_feedback": "LLM not available"}
    llm, _ = llm_tuple

    prompt = ChatPromptTemplate.from_messages([
        ("system", """Review the generated SRS sections for completeness, consistency, and quality. Provide validation feedback and suggest refinements. Output in markdown under 'Validation Feedback'."""),
        MessagesPlaceholder("messages"),
        ("human", "{all_sections}")
    ])
    parser = StrOutputParser()

    all_sections = f"""
Analysis: {state['requirements_analysis']}
Functional: {state['functional_requirements']}
Non-Functional: {state['non_functional_requirements']}
Architecture: {state['system_architecture']}
    """
    feedback = await invoke_llm_with_retry(llm, prompt, parser, {
        "all_sections": all_sections,
        "messages": state["messages"]
    })
    return {
        "validation_feedback": feedback,
        "messages": state["messages"] + [HumanMessage(content=feedback)]
    }

async def compile_final_srs_node(state: AgenticSRSState) -> Dict[str, Any]:
    llm_tuple = get_next_llm()
    if not llm_tuple:
        return {"final_srs": "LLM not available"}
    llm, _ = llm_tuple

    prompt = ChatPromptTemplate.from_messages([
        ("system", """Compile all sections into a professional SRS document in markdown format using the IEEE SRS structure."""),
        MessagesPlaceholder("messages"),
        ("human", "{all_content}")
    ])
    parser = StrOutputParser()

    all_content = f"""
User Input: {state['user_input']}
Analysis: {state['requirements_analysis']}
Functional: {state['functional_requirements']}
Non-Functional: {state['non_functional_requirements']}
Architecture: {state['system_architecture']}
Feedback: {state['validation_feedback']}
    """
    final_srs = await invoke_llm_with_retry(llm, prompt, parser, {
        "all_content": all_content,
        "messages": state["messages"]
    })
    return {
        "final_srs": final_srs,
        "messages": state["messages"] + [HumanMessage(content=final_srs)]
    }

async def ai_edit_content_node(state: AgenticSRSState) -> Dict[str, Any]:
    llm_tuple = get_next_llm()
    if not llm_tuple:
        return {"final_srs": "LLM not available"}
    llm, _ = llm_tuple

    prompt = ChatPromptTemplate.from_messages([
        ("system", """You are an expert document editor with surgical precision. Your task is to make ONLY the specific changes requested by the user while preserving everything else EXACTLY as it is.
        
        CRITICAL RULES:
        1. Read the current content carefully and understand its structure
        2. Identify ONLY the sections that need modification based on the user's instructions
        3. Make minimal, targeted changes - do not rewrite sections unnecessarily
        4. Preserve ALL formatting, markdown syntax, headers, lists, and structure
        5. Maintain the exact same document length and organization unless specifically asked to add/remove content
        6. Do not add explanations, comments, or additional text not requested
        7. Output ONLY the complete document with the requested changes applied
        8. If the instruction is unclear, make the most conservative change possible
        
        Focus on surgical precision - change only what's explicitly requested."""),
        MessagesPlaceholder("messages"),
        ("human", "CURRENT DOCUMENT CONTENT:\n{current_content}\n\n=== EDIT INSTRUCTIONS ===\n{edit_instructions}\n\n=== IMPORTANT ===\nOutput the complete document with ONLY the changes specified in the edit instructions. Preserve all other content exactly as it is.\n\nEDITED DOCUMENT:")
    ])
    parser = StrOutputParser()

    # Add content validation
    current_content = state.get("current_content", "")
    edit_instructions = state.get("edit_instructions", "")
    
    if not current_content.strip():
        return {"final_srs": "No content provided for editing"}
    
    if not edit_instructions.strip():
        return {"final_srs": current_content}  # Return unchanged if no instructions

    edited_content = await invoke_llm_with_retry(llm, prompt, parser, {
        "current_content": current_content,
        "edit_instructions": edit_instructions,
        "messages": state["messages"]
    })
    
    # Validate the edited content
    if not edited_content or len(edited_content.strip()) < len(current_content.strip()) * 0.5:
        logger.warning("AI edit produced significantly shorter content, returning original")
        return {"final_srs": current_content}
    
    return {
        "final_srs": edited_content,
        "messages": state["messages"] + [HumanMessage(content=f"Applied edit instructions: {edit_instructions[:100]}...")]
    }

def build_agentic_srs_workflow():
    workflow = StateGraph(AgenticSRSState)
    workflow.add_node("analyze_requirements", analyze_requirements_node)
    workflow.add_node("generate_functional", generate_functional_requirements_node)
    workflow.add_node("generate_non_functional", generate_non_functional_requirements_node)
    workflow.add_node("generate_architecture", generate_system_architecture_node)
    workflow.add_node("validate_refine", validate_and_refine_node)
    workflow.add_node("compile_final", compile_final_srs_node)

    workflow.set_entry_point("analyze_requirements")
    workflow.add_edge("analyze_requirements", "generate_functional")
    workflow.add_edge("generate_functional", "generate_non_functional")
    workflow.add_edge("generate_non_functional", "generate_architecture")
    workflow.add_edge("generate_architecture", "validate_refine")
    workflow.add_edge("validate_refine", "compile_final")
    workflow.add_edge("compile_final", END)

    checkpointer = MemorySaver()
    return workflow.compile(checkpointer=checkpointer)

agentic_workflow = build_agentic_srs_workflow()

app = FastAPI(
    title="Agentic AI SRS Generator",
    description="Agentic AI-powered Software Requirements Specification generation system with Groq LLM and LangGraph",
    version="2.2.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

sio = socketio.AsyncServer(cors_allowed_origins="*", async_mode="asgi")
socket_app = socketio.ASGIApp(sio, app)

class ChatMessage(BaseModel):
    content: str
    user_input: Optional[str] = None
    document_id: Optional[str] = None

class DocumentResponse(BaseModel):
    id: str
    filename: str
    content: str
    created_at: datetime
    file_type: str

class ExportRequest(BaseModel):
    document_id: str
    format: str
    content: str

class DirectExportRequest(BaseModel):
    content: str
    format: str
    filename: Optional[str] = "srs_document"

class AIEditRequest(BaseModel):
    current_content: str
    edit_instructions: str
    document_id: Optional[str] = None

async def generate_agentic_srs(user_content: str, config: Optional[Dict] = None) -> str:
    if not llms:
        return "LLM not available"
    try:
        initial_state = {
            "messages": [SystemMessage(content="You are an agentic SRS generator.")],
            "user_input": user_content
        }
        result = await agentic_workflow.ainvoke(initial_state, config=config)
        return result.get("final_srs", "No final SRS generated")
    except Exception as e:
        logger.error(f"Error in agentic SRS generation: {e}")
        return f"Error generating SRS: {str(e)}"

async def edit_content_with_ai(current_content: str, edit_instructions: str, config: Optional[Dict] = None) -> str:
    """Enhanced AI content editing with validation"""
    if not llms:
        return "LLM not available"
    
    # Input validation
    if not current_content.strip():
        return "Error: No content provided for editing"
        
    if not edit_instructions.strip():
        return "Error: No edit instructions provided"
    
    # Validate content length
    content_tokens = estimate_tokens(current_content)
    if content_tokens > TOKEN_LIMIT - 1000:
        return "Error: Content too long for editing. Please break into smaller sections."
    
    try:
        # Enhanced validation and processing
        initial_state = {
            "messages": [SystemMessage(content="You are an expert document editor with validation capabilities.")],
            "current_content": current_content,
            "edit_instructions": edit_instructions
        }
        
        logger.info(f"Processing AI edit request - Content length: {len(current_content)} chars")
        result = await ai_edit_content_node(initial_state)
        edited_content = result.get("final_srs", "")
        
        # Post-processing validation
        if not edited_content or len(edited_content.strip()) < len(current_content.strip()) * 0.5:
            logger.warning("AI edit produced significantly shorter content, returning original")
            return current_content  # Return original if edit seems problematic
        
        logger.info(f"AI edit completed successfully - Output length: {len(edited_content)} chars")
        return edited_content
        
    except Exception as e:
        logger.error(f"Error in AI content editing: {e}")
        return f"Error editing content: {str(e)}"

async def stream_ai_edit(current_content: str, edit_instructions: str, config: Optional[Dict] = None) -> AsyncGenerator[str, None]:
    """Enhanced AI edit streaming with proper validation and rate limiting"""
    
    # Input validation
    if not llms:
        yield "LLM not available"
        return
        
    if not current_content.strip():
        yield "Error: No content provided for editing"
        return
        
    if not edit_instructions.strip():
        yield "Error: No edit instructions provided"
        return
    
    # Validate content length and complexity
    content_tokens = estimate_tokens(current_content)
    if content_tokens > TOKEN_LIMIT - 1000:  # Leave room for instructions
        yield "Error: Content too long for editing. Please break into smaller sections."
        return
    
    try:
        llm_tuple = get_next_llm()
        if not llm_tuple:
            yield "No LLM available"
            return
        llm, api_key = llm_tuple
        
        # Enhanced system prompt for better precision
        prompt = ChatPromptTemplate.from_messages([
            ("system", """You are an expert document editor with surgical precision. Your task is to make ONLY the specific changes requested while preserving everything else EXACTLY as it is.

            CRITICAL VALIDATION RULES:
            1. First, read and understand the COMPLETE current content
            2. Identify ONLY the specific sections mentioned in the edit instructions
            3. Validate that the requested changes make sense in the document context
            4. Make minimal, targeted changes - do not rewrite entire sections unless explicitly requested
            5. Preserve ALL formatting, markdown syntax, headers, lists, and structure
            6. Maintain document length and organization unless specifically asked to add/remove content
            7. Do not add explanations, comments, or additional text not requested
            8. Stream the complete document with ONLY the requested changes applied
            9. If instructions are unclear or contradictory, make the most conservative change possible
            
            Focus on surgical precision - change only what's explicitly requested."""),
            MessagesPlaceholder("messages"),
            ("human", """CURRENT DOCUMENT CONTENT:
{current_content}

=== EDIT INSTRUCTIONS ===
{edit_instructions}

=== VALIDATION CHECKLIST ===
Before making changes, ensure:
- You understand the current document structure
- The edit instructions are clear and specific
- Your changes will improve the document as requested
- You preserve all other content exactly as it is

Stream the complete edited document with only the specified changes applied:""")
        ])
        parser = StrOutputParser()

        initial_state = {
            "messages": [SystemMessage(content="You are an expert document editor with validation capabilities.")],
            "current_content": current_content,
            "edit_instructions": edit_instructions
        }

        logger.info(f"Starting AI edit with API key {api_key[-4:]} - Content length: {len(current_content)} chars")
        
        # Stream with enhanced rate limiting and pauses
        full_response = ""
        chunk_count = 0
        
        async for chunk in invoke_llm_stream_with_retry(llm, prompt, parser, {
            "current_content": current_content,
            "edit_instructions": edit_instructions,
            "messages": initial_state["messages"]
        }):
            if chunk:
                full_response += chunk
                chunk_count += 1
                
                # Enhanced streaming with strategic pauses
                yield chunk
                
                # Strategic pauses to reduce API rate limit pressure
                if chunk_count % 10 == 0:  # Every 10 chunks
                    await asyncio.sleep(0.3)  # Longer pause
                elif chunk_count % 5 == 0:  # Every 5 chunks
                    await asyncio.sleep(0.2)  # Medium pause
                else:
                    await asyncio.sleep(0.15)  # Base streaming pace
                    
                # Extra pause at sentence boundaries for better readability
                if chunk.endswith(('.', '!', '?', '\n')):
                    await asyncio.sleep(0.1)
        
        # Validate the edited content
        if full_response and len(full_response.strip()) >= len(current_content.strip()) * 0.5:
            logger.info(f"AI edit completed successfully - Output length: {len(full_response)} chars")
        else:
            logger.warning("AI edit produced significantly shorter content")
            yield "\n\n⚠️ **Warning**: The edited content appears to be significantly shorter than expected. Please review carefully."
            
    except Exception as e:
        logger.error(f"Error streaming AI edit: {e}")
        yield f"\n\n❌ **Error**: {str(e)}"

async def stream_agentic_srs(user_content: str, config: Optional[Dict] = None) -> AsyncGenerator[str, None]:
    if not llms:
        yield "LLM not available"
        return
    try:
        initial_state = {
            "messages": [SystemMessage(content="You are an agentic SRS generator.")],
            "user_input": user_content
        }

        nodes = [
            ("analyze_requirements", analyze_requirements_node, "Analyzing Requirements"),
            ("generate_functional", generate_functional_requirements_node, "Generating Functional Requirements"),
            ("generate_non_functional", generate_non_functional_requirements_node, "Generating Non-Functional Requirements"),
            ("generate_architecture", generate_system_architecture_node, "Creating System Architecture"),
            ("validate_refine", validate_and_refine_node, "Validating and Refining"),
            ("compile_final", compile_final_srs_node, "Compiling Final SRS")
        ]
        state = initial_state

        for node, node_func, step_name in nodes:
            partial = await node_func(state)
            state.update(partial)
            key = next(iter(partial))
            content = partial.get(key, "")
            if content:
                # Enhanced formatting for better display
                formatted = format_streaming_content(content)
                # Stream word by word with proper spacing
                words = formatted.split()
                current_line = ""
                for i, word in enumerate(words):
                    current_line += word + " "
                    # Stream complete lines for better readability
                    if word.endswith(('.', '!', '?', ':')) or '\n' in word or i == len(words) - 1:
                        yield current_line.strip() + "\n"
                        current_line = ""
                        await asyncio.sleep(0.08)  # Balanced streaming speed
                    elif len(current_line) > 80:  # Stream long lines in chunks
                        yield current_line.strip() + " "
                        current_line = ""
                        await asyncio.sleep(0.05)
            yield f"\n\n✓ **{step_name} Completed**\n\n"
            await asyncio.sleep(0.3)
        yield "# Agentic SRS generation complete.\n\n"
    except Exception as e:
        logger.error(f"Error streaming agentic SRS: {e}")
        yield f"Error streaming SRS: {str(e)}"

def get_file_type(file_path: Path) -> str:
    if MAGIC_AVAILABLE:
        try:
            return magic.from_file(str(file_path), mime=True)
        except Exception as e:
            logger.warning(f"python-magic failed: {e}. Falling back to extension.")
    
    extension = file_path.suffix.lower()
    mime_types = {
        '.txt': 'text/plain',
        '.pdf': 'application/pdf',
        '.doc': 'application/msword',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.ppt': 'application/vnd.ms-powerpoint',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.jpg': 'image/jpeg',
        '.jpeg': 'image/jpeg',
        '.png': 'image/png',
        '.gif': 'image/gif',
        '.json': 'application/json',
    }
    return mime_types.get(extension, 'application/octet-stream')

async def extract_text_from_file(file_path: Path, file_type: str) -> str:
    try:
        if file_type.startswith('text/') or file_type == 'application/json':
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                return await f.read()
        elif file_type == 'application/pdf':
            return extract_pdf_text(file_path)
        elif file_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
            return extract_docx_text(file_path)
        elif file_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
            return extract_pptx_text(file_path)
        elif file_type.startswith('image/') and ocr:
            return extract_text_from_image(file_path)
        else:
            async with aiofiles.open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return await f.read()
    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        return ""

def extract_pdf_text(file_path: Path) -> str:
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
            return text
    except Exception as e:
        logger.error(f"Error extracting PDF text: {e}")
        return ""

def extract_docx_text(file_path: Path) -> str:
    try:
        doc = docx.Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting DOCX text: {e}")
        return ""

def extract_pptx_text(file_path: Path) -> str:
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PPTX text: {e}")
        return ""

def extract_text_from_image(image_path: Path) -> str:
    if not ocr:
        return ""
    try:
        results = ocr.readtext(str(image_path))
        text = ""
        for _, detected_text, conf in results:
            if conf > 0.5:
                text += detected_text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting text from image: {e}")
        return ""

def clean_markdown_formatting(content: str) -> str:
    content = re.sub(r'^#{1,6}\s+', '', content, flags=re.MULTILINE)
    content = re.sub(r'\*\*(.*?)\*\*', r'\1', content)
    content = re.sub(r'\*(.*?)\*', r'\1', content)
    content = re.sub(r'__(.*?)__', r'\1', content)
    content = re.sub(r'_(.*?)_', r'\1', content)
    content = re.sub(r'^[\-\*\+]\s+', '', content, flags=re.MULTILINE)
    content = re.sub(r'^\d+\.\s+', '', content, flags=re.MULTILINE)
    content = re.sub(r'```[\s\S]*?```', '', content)
    content = re.sub(r'`([^`]+)`', r'\1', content)
    content = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', content)
    content = re.sub(r'^[-*_]{3,}$', '', content, flags=re.MULTILINE)
    content = re.sub(r'^>\s+', '', content, flags=re.MULTILINE)
    return content.strip()

def is_title_line(line: str) -> bool:
    title_indicators = ['software requirements specification', 'srs document', 'system requirements', 'project requirements']
    return any(indicator in line.lower() for indicator in title_indicators) and len(line) < 100

def is_main_heading(line: str) -> bool:
    heading_patterns = [
        r'^\d+\.',  # Numbered sections like "1. Introduction"
        r'^[A-Z][A-Z\s]{3,}$',  # ALL CAPS headings
    ]
    main_headings = ['introduction', 'overview', 'requirements', 'architecture', 'design', 'implementation', 'testing', 'conclusion']
    if any(re.match(pattern, line) for pattern in heading_patterns):
        return True
    return any(heading in line.lower() for heading in main_headings) and len(line.split()) <= 5

def is_sub_heading(line: str) -> bool:
    sub_patterns = [
        r'^\d+\.\d+',  # Numbered subsections like "1.1 Purpose"
        r'^[A-Z][a-z]+\s[A-Z][a-z]+',  # Title Case headings
    ]
    return any(re.match(pattern, line) for pattern in sub_patterns) or (
        line.istitle() and len(line.split()) <= 6 and len(line) < 80
    )

def create_pdf_export(content: str, output_path: Path) -> Path:
    try:
        doc = SimpleDocTemplate(
            str(output_path), 
            pagesize=A4,
            rightMargin=72, 
            leftMargin=72,
            topMargin=72, 
            bottomMargin=72
        )
        styles = getSampleStyleSheet()
        
        # Enhanced styles with better formatting
        title_style = ParagraphStyle(
            'CustomTitle', 
            parent=styles['Title'], 
            fontSize=18, 
            spaceAfter=30, 
            alignment=TA_CENTER,
            textColor='#2C3E50',
            fontName='Helvetica-Bold'
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading', 
            parent=styles['Heading1'], 
            fontSize=14, 
            spaceAfter=12, 
            spaceBefore=18,
            textColor='#34495E',
            fontName='Helvetica-Bold'
        )
        
        subheading_style = ParagraphStyle(
            'CustomSubHeading', 
            parent=styles['Heading2'], 
            fontSize=12, 
            spaceAfter=8, 
            spaceBefore=12,
            textColor='#5D6D7E',
            fontName='Helvetica-Bold'
        )
        
        body_style = ParagraphStyle(
            'CustomBody', 
            parent=styles['Normal'], 
            fontSize=10, 
            spaceAfter=6, 
            alignment=TA_JUSTIFY,
            fontName='Helvetica'
        )
        
        list_style = ParagraphStyle(
            'CustomList', 
            parent=styles['Normal'], 
            fontSize=10, 
            spaceAfter=3,
            leftIndent=20,
            fontName='Helvetica'
        )
        
        flowables = []
        lines = content.split('\n')
        
        # Add title if content starts with a heading
        first_heading = next((line.strip() for line in lines if line.strip().startswith('#')), None)
        title_text = "Software Requirements Specification"  # Default title
        if first_heading:
            title_text = re.sub(r'^#+\s*', '', first_heading)
            flowables.append(Paragraph(title_text, title_style))
            flowables.append(Spacer(1, 20))
        else:
            flowables.append(Paragraph(title_text, title_style))
            flowables.append(Spacer(1, 20))
        
        for line in lines:
            line = line.strip()
            if not line:
                flowables.append(Spacer(1, 6))
                continue
                
            # Handle markdown headings
            if line.startswith('# '):
                text = line[2:].strip()
                if text != title_text:  # Don't repeat the title
                    flowables.append(Spacer(1, 12))
                    flowables.append(Paragraph(escape_html(text), heading_style))
            elif line.startswith('## '):
                text = line[3:].strip()
                flowables.append(Spacer(1, 8))
                flowables.append(Paragraph(escape_html(text), subheading_style))
            elif line.startswith('### '):
                text = line[4:].strip()
                flowables.append(Paragraph(escape_html(text), subheading_style))
            elif re.match(r'^\d+\.\s', line):  # Numbered list
                text = escape_html(line)
                flowables.append(Paragraph(text, list_style))
            elif re.match(r'^[-*+]\s', line):  # Bullet list
                text = escape_html(line)
                flowables.append(Paragraph(text, list_style))
            else:
                # Regular paragraph
                if line and not line.startswith('#'):
                    # Handle bold and italic markdown
                    text = process_markdown_formatting(escape_html(line))
                    flowables.append(Paragraph(text, body_style))
        
        doc.build(flowables)
        logger.info(f"PDF created successfully: {output_path}")
        return output_path
    except Exception as e:
        logger.error(f"Error creating PDF: {e}")
        raise

def escape_html(text: str) -> str:
    """Escape HTML characters for reportlab."""
    return text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

def process_markdown_formatting(text: str) -> str:
    """Convert basic markdown formatting to HTML for reportlab."""
    # Bold text
    text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text)
    # Italic text
    text = re.sub(r'\*(.*?)\*', r'<i>\1</i>', text)
    # Inline code
    text = re.sub(r'`(.*?)`', r'<font name="Courier">\1</font>', text)
    return text

def create_docx_export(content: str, output_path: Path) -> Path:
    try:
        doc = docx.Document()
        
        # Set document margins
        sections = doc.sections
        for section in sections:
            section.top_margin = docx.shared.Inches(1)
            section.bottom_margin = docx.shared.Inches(1)
            section.left_margin = docx.shared.Inches(1)
            section.right_margin = docx.shared.Inches(1)
        
        lines = content.split('\n')
        
        # Add title from first heading
        first_heading = next((line.strip() for line in lines if line.strip().startswith('#')), None)
        if first_heading:
            title_text = re.sub(r'^#+\s*', '', first_heading)
            title = doc.add_heading(title_text, 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            doc.add_paragraph()  # Add space after title
        
        for line in lines:
            line = line.strip()
            if not line:
                doc.add_paragraph()  # Empty paragraph for spacing
                continue
                
            # Handle markdown headings
            if line.startswith('# '):
                text = line[2:].strip()
                if first_heading and text == re.sub(r'^#+\s*', '', first_heading):
                    continue  # Skip duplicate title
                heading = doc.add_heading(text, level=1)
                heading.style = 'Heading 1'
            elif line.startswith('## '):
                text = line[3:].strip()
                heading = doc.add_heading(text, level=2)
                heading.style = 'Heading 2'
            elif line.startswith('### '):
                text = line[4:].strip()
                heading = doc.add_heading(text, level=3)
                heading.style = 'Heading 3'
            elif re.match(r'^\d+\.\s', line):  # Numbered list
                p = doc.add_paragraph(line, style='List Number')
            elif re.match(r'^[-*+]\s', line):  # Bullet list
                # Remove the markdown bullet and add as Word bullet
                text = re.sub(r'^[-*+]\s', '', line)
                p = doc.add_paragraph(text, style='List Bullet')
            else:
                # Regular paragraph with markdown formatting
                if line and not line.startswith('#'):
                    p = doc.add_paragraph()
                    add_formatted_text(p, line)
        
        doc.save(str(output_path))
        logger.info(f"DOCX created successfully: {output_path}")
        return output_path
    except Exception as e:
        logger.error(f"Error creating DOCX: {e}")
        raise

def add_formatted_text(paragraph, text: str):
    """Add text with markdown formatting to a Word paragraph."""
    # Split text by markdown patterns and apply formatting
    parts = re.split(r'(\*\*.*?\*\*|\*.*?\*|`.*?`)', text)
    
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            # Bold text
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        elif part.startswith('*') and part.endswith('*') and not part.startswith('**'):
            # Italic text
            run = paragraph.add_run(part[1:-1])
            run.italic = True
        elif part.startswith('`') and part.endswith('`'):
            # Code text
            run = paragraph.add_run(part[1:-1])
            run.font.name = 'Courier New'
        else:
            # Regular text
            if part:  # Only add non-empty parts
                paragraph.add_run(part)


@app.get("/")
async def home():
    return {
        "message": "Welcome to Agentic AI SRS Generator with Groq LLM",
        "version": "2.2.0",
        "status": "active",
        "llm_available": len(llms) > 0,
        "magic_available": MAGIC_AVAILABLE,
        "agentic_workflow": "active"
    }

@app.get("/health")
async def health_check():
    return {
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "ocr_available": ocr is not None,
        "llm_available": len(llms) > 0,
        "magic_available": MAGIC_AVAILABLE,
        "num_llms": len(llms),
        "workflow_nodes": 6
    }

class SessionMemoryManager:
    @staticmethod
    def create_session(session_id: str, user_input: str, document_id: Optional[str] = None, uploaded_files: Optional[List] = None) -> str:
        relevant_knowledge = []
        context_summary = ""
        
        # Process uploaded files for contextual knowledge
        if uploaded_files:
            for file_info in uploaded_files:
                doc_id = file_info.get("document_id")
                if doc_id in document_knowledge_base:
                    knowledge_item = document_knowledge_base[doc_id]
                    relevant_knowledge.append(knowledge_item)
        
        # Get relevant knowledge based on user input (for context only)
        query_knowledge = DocumentLearningAgent.get_relevant_knowledge(user_input, limit=3)
        for knowledge_item in query_knowledge:
            # Avoid duplicates
            knowledge_data = knowledge_item["knowledge"]
            if not any(existing.get("document_id") == knowledge_data.get("document_id") for existing in relevant_knowledge):
                relevant_knowledge.append(knowledge_data)
        
        # Create context summary emphasizing it's for reference only
        if relevant_knowledge:
            context_parts = []
            for knowledge in relevant_knowledge[:3]:  # Limit to top 3
                if isinstance(knowledge, dict):
                    filename = knowledge.get('filename', 'Unknown Document')
                    analysis = knowledge.get('analysis', '')
                    # Extract key contextual insights
                    domain_context = analysis[:400] + "..." if len(analysis) > 400 else analysis
                    context_parts.append(f"Context from {filename}: {domain_context}")
            context_summary = "\n\n".join(context_parts)
        
        initial_system_message = SystemMessage(content="""You are an agentic SRS generator. You will analyze USER REQUIREMENTS to create comprehensive SRS documents. 
        
        IMPORTANT: Any document context provided is REFERENCE MATERIAL ONLY to help you understand:
        - Domain terminology and standards
        - Industry best practices
        - Technical architecture patterns
        - Quality attributes and compliance considerations
        
        Always focus on the USER'S INPUT REQUIREMENTS as the primary source for SRS generation.""")
        
        session_workflows[session_id] = {
            "id": session_id,
            "user_input": user_input,
            "document_id": document_id,
            "uploaded_files": uploaded_files or [],
            "relevant_knowledge": relevant_knowledge,
            "context_summary": context_summary,
            "created_at": datetime.now(),
            "current_step": "analyze_requirements",
            "completed_steps": [],
            "state": {
                "messages": [initial_system_message],
                "user_input": user_input,
                "document_context": relevant_knowledge,  # Structured context data
                "requirements_analysis": None,
                "functional_requirements": None,
                "non_functional_requirements": None,
                "system_architecture": None,
                "validation_feedback": None,
                "final_srs": None,
                "edit_context": None,
                "edit_instructions": None,
                "current_content": None
            }
        }
        
        session_progress[session_id] = {
            "total_steps": 6,
            "completed_steps": 0,
            "current_step_name": "Analyzing Requirements",
            "progress_percentage": 0,
            "status": "in_progress",
            "start_time": datetime.now(),
            "estimated_completion": None,
            "has_document_context": len(relevant_knowledge) > 0,
            "document_count": len(uploaded_files) if uploaded_files else 0,
            "step_details": {
                "analyze_requirements": {"status": "pending", "result_summary": None, "timestamp": None},
                "generate_functional": {"status": "pending", "result_summary": None, "timestamp": None},
                "generate_non_functional": {"status": "pending", "result_summary": None, "timestamp": None},
                "generate_architecture": {"status": "pending", "result_summary": None, "timestamp": None},
                "validate_refine": {"status": "pending", "result_summary": None, "timestamp": None},
                "compile_final": {"status": "pending", "result_summary": None, "timestamp": None}
            }
        }
        
        logger.info(f"Created session {session_id} with {len(relevant_knowledge)} relevant documents")
        return session_id
    
    @staticmethod
    def update_session_step(session_id: str, step_name: str, result: Any) -> None:
        if session_id not in session_workflows:
            return
        session_workflows[session_id]["current_step"] = step_name
        session_workflows[session_id]["completed_steps"].append(step_name)
        session_workflows[session_id]["updated_at"] = datetime.now()
        
        if session_id in session_progress:
            progress = session_progress[session_id]
            progress["completed_steps"] += 1
            progress["progress_percentage"] = (progress["completed_steps"] / progress["total_steps"]) * 100
            
            step_names = {
                "analyze_requirements": "Analyzing Requirements",
                "generate_functional": "Generating Functional Requirements",
                "generate_non_functional": "Generating Non-Functional Requirements",
                "generate_architecture": "Creating System Architecture",
                "validate_refine": "Validating and Refining",
                "compile_final": "Compiling Final SRS"
            }
            
            progress["current_step_name"] = step_names.get(step_name, "Processing")
            progress["step_details"][step_name] = {
                "status": "completed",
                "result_summary": str(result)[:500] if result else None,
                "timestamp": datetime.now()
            }
            
            if progress["completed_steps"] >= progress["total_steps"]:
                progress["status"] = "completed"
                progress["progress_percentage"] = 100
                progress["completion_time"] = datetime.now()
            
            logger.info(f"Session {session_id}: Completed {step_name} ({progress['progress_percentage']:.1f}%)")
    
    @staticmethod
    def get_session_progress(session_id: str) -> Optional[Dict[str, Any]]:
        progress = session_progress.get(session_id)
        if not progress:
            return None
        return safe_json_serialize(progress)
    
    @staticmethod
    def get_session_state(session_id: str) -> Optional[Dict[str, Any]]:
        workflow = session_workflows.get(session_id)
        if not workflow:
            return None
        return copy.deepcopy(workflow["state"])
    
    @staticmethod
    def get_session_context(session_id: str) -> str:
        workflow = session_workflows.get(session_id)
        if not workflow or not workflow.get("relevant_knowledge"):
            return ""
        context_parts = []
        for knowledge in workflow["relevant_knowledge"]:
            context_parts.append(f"\nDocument: {knowledge.get('filename', 'Unknown')}")
            context_parts.append(f"Analysis: {knowledge.get('analysis', 'No analysis available')}")
        return "\n".join(context_parts)
    
    @staticmethod
    def update_session_state(session_id: str, key: str, value: Any) -> None:
        if session_id in session_workflows:
            session_workflows[session_id]["state"][key] = value
    
    @staticmethod
    def cleanup_old_sessions(max_age_hours: int = 24) -> None:
        cutoff_time = datetime.now() - timedelta(hours=max_age_hours)
        sessions_to_remove = []
        for session_id, workflow in session_workflows.items():
            if workflow.get("created_at", datetime.now()) < cutoff_time:
                sessions_to_remove.append(session_id)
        
        for session_id in sessions_to_remove:
            session_workflows.pop(session_id, None)
            session_progress.pop(session_id, None)
            logger.info(f"Cleaned up old session: {session_id}")
    
    @staticmethod
    def get_all_sessions() -> Dict[str, Any]:
        sessions_info = {
            "active_sessions": len(session_workflows),
            "sessions": {}
        }
        for session_id, workflow in session_workflows.items():
            progress = SessionMemoryManager.get_session_progress(session_id)
            sessions_info["sessions"][session_id] = {
                "created_at": workflow["created_at"].isoformat(),
                "progress": progress,
                "document_id": workflow.get("document_id")
            }
        return sessions_info

@app.websocket("/stream-srs")
async def websocket_stream_srs(websocket: WebSocket):
    await websocket.accept()
    connection_id = str(uuid.uuid4())
    websocket_connections[connection_id] = websocket
    logger.info(f"WebSocket connection established: {connection_id}")
    
    try:
        while True:
            data = await websocket.receive_text()
            message_data = json.loads(data)
            
            content = message_data.get("content", "")
            document_id = message_data.get("document_id")
            uploaded_files = message_data.get("uploaded_files", [])
            
            logger.info(f"Received streaming request: {content[:100]}...")
            
            if not content.strip():
                await websocket.send_text(json.dumps({
                    "error": "No content provided"
                }, default=safe_json_serialize))
                continue
            
            session_id = str(uuid.uuid4())
            SessionMemoryManager.create_session(session_id, content, document_id, uploaded_files)
            
            progress = SessionMemoryManager.get_session_progress(session_id)
            await websocket.send_text(json.dumps({
                "session_created": session_id,
                "progress": progress
            }, default=safe_json_serialize))
            
            state = SessionMemoryManager.get_session_state(session_id)
            full_response = ""
            nodes = [
                ("analyze_requirements", analyze_requirements_node, "Analyzing Requirements"),
                ("generate_functional", generate_functional_requirements_node, "Generating Functional Requirements"),
                ("generate_non_functional", generate_non_functional_requirements_node, "Generating Non-Functional Requirements"),
                ("generate_architecture", generate_system_architecture_node, "Creating System Architecture"),
                ("validate_refine", validate_and_refine_node, "Validating and Refining"),
                ("compile_final", compile_final_srs_node, "Compiling Final SRS")
            ]
            
            for i, (node, node_func, step_name) in enumerate(nodes):
                if session_id in session_progress:
                    session_progress[session_id]["current_step_name"] = step_name
                progress = SessionMemoryManager.get_session_progress(session_id)
                await websocket.send_text(json.dumps({
                    "progress_update": progress
                }, default=safe_json_serialize))
                
                partial = await node_func(state)
                state.update(partial)
                key = next(iter(partial))
                chunk_content = partial.get(key, "")
                SessionMemoryManager.update_session_step(session_id, node, chunk_content)
                session_workflows[session_id]["state"] = state
                
                if chunk_content:
                    full_response += chunk_content + "\n\n"
                    # Enhanced formatting for streaming
                    formatted = format_streaming_content(chunk_content)
                    
                    # Stream in sentence-based chunks for better readability
                    sentences = re.split(r'(?<=[.!?])\s+', formatted)
                    for sentence in sentences:
                        if sentence.strip():
                            await websocket.send_text(json.dumps({
                                "chunk": sentence.strip() + " ",
                                "step": node,
                                "progress": SessionMemoryManager.get_session_progress(session_id)
                            }, default=safe_json_serialize))
                            await asyncio.sleep(0.2)  # Balanced streaming speed
                    
                    await websocket.send_text(json.dumps({
                        "chunk": f"\n\n✓ **{step_name} Completed**\n\n",
                        "step": node,
                        "progress": SessionMemoryManager.get_session_progress(session_id)
                    }, default=safe_json_serialize))
                
                await asyncio.sleep(0.5)
            
            final_progress = SessionMemoryManager.get_session_progress(session_id)
            await websocket.send_text(json.dumps({
                "final": {
                    "final_srs": full_response,
                    "session_id": session_id
                },
                "progress": final_progress
            }, default=safe_json_serialize))
            
            logger.info(f"Streaming agentic SRS generation completed for session: {session_id}")
    
    except WebSocketDisconnect:
        logger.info(f"WebSocket disconnected: {connection_id}")
        websocket_connections.pop(connection_id, None)
    except Exception as e:
        logger.error(f"WebSocket error: {e}")
        await websocket.send_text(json.dumps({
            "error": f"WebSocket error: {str(e)}"
        }, default=safe_json_serialize))
        websocket_connections.pop(connection_id, None)

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")
    
    document_id = str(uuid.uuid4())
    file_extension = Path(file.filename).suffix
    file_path = UPLOAD_DIR / f"{document_id}{file_extension}"
    
    try:
        async with aiofiles.open(file_path, 'wb') as f:
            content = await file.read()
            await f.write(content)
        
        file_type = get_file_type(file_path)
        extracted_text = await extract_text_from_file(file_path, file_type)
        
        documents_store[document_id] = {
            "id": document_id,
            "filename": file.filename,
            "file_path": str(file_path),
            "file_type": file_type,
            "content": extracted_text,
            "created_at": datetime.now(),
            "size": len(content)
        }
        
        try:
            analysis = await DocumentLearningAgent.analyze_document(extracted_text, documents_store[document_id])
            if "error" not in analysis:
                DocumentLearningAgent.store_document_knowledge(document_id, analysis)
                logger.info(f"Document {document_id} analyzed and stored")
        except Exception as e:
            logger.warning(f"Failed to analyze document {document_id}: {e}")
        
        return {
            "status": "success",
            "document_id": document_id,
            "filename": file.filename,
            "file_type": file_type,
            "content_preview": extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text,
            "size": len(content)
        }
    except Exception as e:
        logger.error(f"Error uploading file: {e}")
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.post("/chat")
async def chat_with_ai(message: ChatMessage):
    if not llms:
        raise HTTPException(status_code=503, detail="LLM not available")
    try:
        user_content = message.content
        document_context = ""
        if message.document_id and message.document_id in documents_store:
            doc = documents_store[message.document_id]
            document_context = f"\n\nDocument Content:\n{doc['content']}"
        
        full_input = user_content + document_context
        logger.info("Starting agentic SRS generation...")
        result = await generate_agentic_srs(full_input)
        logger.info("Agentic SRS generation completed")
        return {
            "response": result,
            "status": "success"
        }
    except RateLimitError as e:
        headers = getattr(e, 'response', {}).headers if hasattr(e, 'response') else {}
        retry_after = headers.get("retry-after", "60")
        raise HTTPException(
            status_code=429,
            detail=f"Rate limit exceeded. Try again in {retry_after} seconds."
        )
    except Exception as e:
        logger.error(f"Error in chat endpoint: {e}")
        raise HTTPException(status_code=500, detail=f"Error generating SRS: {str(e)}")

@app.get("/documents/{document_id}")
async def get_document(document_id: str):
    if document_id not in documents_store:
        raise HTTPException(status_code=404, detail="Document not found")
    
    doc = documents_store[document_id]
    return DocumentResponse(
        id=doc["id"],
        filename=doc["filename"],
        content=doc["content"],
        created_at=doc["created_at"],
        file_type=doc["file_type"]
    )

@app.get("/documents")
async def list_documents():
    docs = []
    for doc in documents_store.values():
        docs.append({
            "id": doc["id"],
            "filename": doc["filename"],
            "file_type": doc["file_type"],
            "created_at": doc["created_at"].isoformat(),
            "size": doc["size"]
        })
    return {"documents": docs}

@app.put("/documents/{document_id}")
async def update_document(document_id: str, content: dict):
    if document_id not in documents_store:
        raise HTTPException(status_code=404, detail="Document not found")
    
    documents_store[document_id]["content"] = content.get("content", "")
    documents_store[document_id]["updated_at"] = datetime.now()
    return {"status": "success", "document_id": document_id}


@app.post("/export/{document_id}")
async def export_document(document_id: str, export_request: ExportRequest):
    if document_id not in documents_store:
        raise HTTPException(status_code=404, detail="Document not found")
    
    # Validate export format
    supported_formats = ['pdf', 'docx']
    if export_request.format.lower() not in supported_formats:
        raise HTTPException(status_code=400, detail=f"Unsupported format. Supported formats: {', '.join(supported_formats)}")
    
    try:
        doc = documents_store[document_id]
        content = export_request.content or doc["content"]
        
        # Validate content
        if not content or not content.strip():
            raise HTTPException(status_code=400, detail="No content to export")
        
        logger.info(f"Exporting document {document_id} to {export_request.format} format")
        export_path = await create_export_file(document_id, content, export_request.format)
        
        # Verify the exported file exists and has content
        if not export_path.exists() or export_path.stat().st_size == 0:
            raise HTTPException(status_code=500, detail="Export file was not created successfully")
        
        # Enhanced media type mapping
        media_types = {
            'pdf': 'application/pdf',
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        }
        media_type = media_types.get(export_request.format.lower(), 'application/octet-stream')
        
        # Generate appropriate filename
        base_filename = Path(doc['filename']).stem if doc.get('filename') else 'srs_document'
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{base_filename}_srs_{timestamp}.{export_request.format}"
        
        logger.info(f"Serving export file: {export_path} as {filename}")
        return FileResponse(
            path=str(export_path),
            filename=filename,
            media_type=media_type,
            headers={
                "Content-Disposition": f"attachment; filename={filename}",
                "Cache-Control": "no-cache"
            }
        )
    except HTTPException:
        raise  # Re-raise HTTP exceptions
    except Exception as e:
        logger.error(f"Error exporting document {document_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Error exporting document: {str(e)}")

@app.post("/export-direct")
async def export_direct_content(export_request: DirectExportRequest):
    """Export content directly without requiring a document ID."""
    
    # Validate export format
    supported_formats = ['pdf', 'docx']
    if export_request.format.lower() not in supported_formats:
        raise HTTPException(status_code=400, detail=f"Unsupported format. Supported formats: {', '.join(supported_formats)}")
    
    # Validate content
    if not export_request.content or not export_request.content.strip():
        raise HTTPException(status_code=400, detail="No content to export")
    
    try:
        # Generate a temporary document ID for file creation
        temp_doc_id = f"temp_{uuid.uuid4().hex[:8]}"
        
        logger.info(f"Exporting direct content to {export_request.format} format")
        export_path = await create_export_file(temp_doc_id, export_request.content, export_request.format)
        
        # Verify the exported file exists and has content
        if not export_path.exists() or export_path.stat().st_size == 0:
            raise HTTPException(status_code=500, detail="Export file was not created successfully")
        
        # Media type mapping
        media_types = {
            'pdf': 'application/pdf',
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        }
        media_type = media_types.get(export_request.format.lower(), 'application/octet-stream')
        
        # Generate filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{export_request.filename}_srs_{timestamp}.{export_request.format}"
        
        logger.info(f"Serving direct export file: {export_path} as {filename}")
        return FileResponse(
            path=str(export_path),
            filename=filename,
            media_type=media_type,
            headers={
                "Content-Disposition": f"attachment; filename={filename}",
                "Cache-Control": "no-cache"
            }
        )
    except HTTPException:
        raise  # Re-raise HTTP exceptions
    except Exception as e:
        logger.error(f"Error exporting direct content: {e}")
        raise HTTPException(status_code=500, detail=f"Error exporting content: {str(e)}")

async def create_export_file(document_id: str, content: str, format: str) -> Path:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    export_filename = f"{document_id}_{timestamp}_export.{format}"
    export_path = OUTPUT_DIR / export_filename
    
    try:
        if format == "pdf":
            return create_pdf_export(content, export_path)
        elif format == "docx":
            return create_docx_export(content, export_path)
        else:
            raise ValueError(f"Unsupported export format: {format}")
    except Exception as e:
        logger.error(f"Error creating export file in {format} format: {e}")
        raise


@sio.event
async def connect(sid, environ):
    logger.info(f"Client {sid} connected")
    await sio.emit('connected', {'status': 'connected'}, room=sid)

@sio.event
async def disconnect(sid):
    logger.info(f"Client {sid} disconnected")
    if sid in active_sessions:
        del active_sessions[sid]

@sio.event
async def join_document(sid, data):
    document_id = data.get('document_id')
    user_id = data.get('user_id', f'user_{sid[:8]}')
    
    if document_id:
        active_sessions[sid] = {
            'user_id': user_id,
            'document_id': document_id
        }
        await sio.enter_room(sid, document_id)
        await sio.emit('user_joined', {
            'user_id': user_id,
            'document_id': document_id
        }, room=document_id, skip_sid=sid)
        
        if document_id in documents_store:
            doc = documents_store[document_id]
            await sio.emit('document_state', {
                'content': doc['content'],
                'document_id': document_id
            }, room=sid)

@sio.event
async def content_change(sid, data):
    if sid in active_sessions:
        session = active_sessions[sid]
        document_id = session['document_id']
        content = data.get('content', '')
        
        if document_id in documents_store:
            documents_store[document_id]['content'] = content
            documents_store[document_id]['updated_at'] = datetime.now()
            await sio.emit('content_updated', {
                'content': content,
                'user_id': session['user_id'],
                'timestamp': datetime.now().isoformat()
            }, room=document_id, skip_sid=sid)

@app.get("/workflow/status")
async def get_workflow_status():
    return {
        "llm_available": len(llms) > 0,
        "ocr_available": ocr is not None,
        "magic_available": MAGIC_AVAILABLE,
        "num_llms": len(llms),
        "supported_formats": ["pdf", "docx"],
        "supported_upload_types": [
            "text/*", "application/pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "image/*"
        ],
        "websocket_connections": len(websocket_connections),
        "agentic_nodes": ["analyze_requirements", "generate_functional", "generate_non_functional", "generate_architecture", "validate_refine", "compile_final"]
    }

@app.post("/workflow/test")
async def test_workflow():
    if not llms:
        raise HTTPException(status_code=503, detail="LLM not available")
    
    test_input = """
    Test Requirements:
    - Create a web-based task management system
    - Users should be able to create, edit, and delete tasks
    - System should support user authentication
    - Tasks should have priorities and due dates
    - System should send email notifications for due tasks
    """
    
    try:
        result = await generate_agentic_srs(test_input)
        return {
            "status": "success",
            "message": "Agentic workflow test completed successfully",
            "preview": result[:500] + "..." if len(result) > 500 else result
        }
    except Exception as e:
        logger.error(f"Workflow test failed: {e}")
        return {
            "status": "error",
            "message": f"Workflow test failed: {str(e)}"
        }

@app.delete("/documents/{document_id}")
async def delete_document(document_id: str):
    if document_id not in documents_store:
        raise HTTPException(status_code=404, detail="Document not found")
    
    try:
        doc = documents_store[document_id]
        file_path = Path(doc["file_path"])
        if file_path.exists():
            file_path.unlink()
        del documents_store[document_id]
        document_knowledge_base.pop(document_id, None)
        return {
            "status": "success",
            "message": f"Document {document_id} deleted successfully"
        }
    except Exception as e:
        logger.error(f"Error deleting document: {e}")
        raise HTTPException(status_code=500, detail=f"Error deleting document: {str(e)}")

@app.get("/sessions")
async def list_sessions():
    return SessionMemoryManager.get_all_sessions()

@app.get("/sessions/{session_id}/progress")
async def get_session_progress(session_id: str):
    progress = SessionMemoryManager.get_session_progress(session_id)
    if not progress:
        raise HTTPException(status_code=404, detail="Session not found")
    return progress

@app.get("/knowledge-base")
async def get_knowledge_base():
    return DocumentLearningAgent.get_knowledge_summary()

@app.get("/knowledge-base/search")
async def search_knowledge_base(query: str, limit: int = 5):
    if not query.strip():
        raise HTTPException(status_code=400, detail="Query cannot be empty")
    
    results = DocumentLearningAgent.get_relevant_knowledge(query, limit)
    return {
        "query": query,
        "results": results,
        "total_found": len(results)
    }

@app.post("/sessions/cleanup")
async def cleanup_old_sessions(max_age_hours: int = 24):
    initial_count = len(session_workflows)
    SessionMemoryManager.cleanup_old_sessions(max_age_hours)
    cleaned_count = initial_count - len(session_workflows)
    return {
        "status": "success",
        "message": f"Cleaned up {cleaned_count} old sessions",
        "remaining_sessions": len(session_workflows)
    }

@app.post("/documents/{document_id}/analyze")
async def analyze_document(document_id: str):
    if document_id not in documents_store:
        raise HTTPException(status_code=404, detail="Document not found")
    
    if not llms:
        raise HTTPException(status_code=503, detail="LLM not available")
    
    try:
        doc = documents_store[document_id]
        content = doc["content"]
        
        analysis_input = f"""
        Analyze the following document content and extract software requirements:
        
        Document: {doc['filename']}
        Content: {content}
        
        Please identify:
        1. Business requirements and objectives
        2. Functional requirements
        3. Non-functional requirements
        4. Constraints and assumptions
        5. Stakeholders and their needs
        
        Provide a structured analysis in markdown format.
        """
        
        result = await generate_agentic_srs(analysis_input)
        
        documents_store[document_id]["analysis"] = {
            "analysis_result": result,
            "analyzed_at": datetime.now()
        }
        
        return {
            "status": "success",
            "document_id": document_id,
            "analysis": result[:1000] + "..." if len(result) > 1000 else result,
            "full_analysis_available": True
        }
    except Exception as e:
        logger.error(f"Error analyzing document: {e}")
        raise HTTPException(status_code=500, detail=f"Error analyzing document: {str(e)}")

@app.post("/ai-edit")
async def ai_edit_content(edit_request: AIEditRequest):
    if not llms:
        raise HTTPException(status_code=503, detail="LLM not available")
    
    try:
        result = await edit_content_with_ai(
            edit_request.current_content, 
            edit_request.edit_instructions
        )
        return {
            "status": "success",
            "edited_content": result
        }
    except Exception as e:
        logger.error(f"Error in AI editing: {e}")
        raise HTTPException(status_code=500, detail=f"Error editing content: {str(e)}")

@app.websocket("/stream-ai-edit")
async def websocket_stream_ai_edit(websocket: WebSocket):
    await websocket.accept()
    connection_id = str(uuid.uuid4())
    websocket_connections[connection_id] = websocket
    logger.info(f"AI Edit WebSocket connection established: {connection_id}")
    
    try:
        while True:
            data = await websocket.receive_text()
            message_data = json.loads(data)
            
            current_content = message_data.get("current_content", "")
            edit_instructions = message_data.get("edit_instructions", "")
            
            logger.info(f"Received AI edit request: {edit_instructions[:100]}...")
            
            if not current_content.strip() or not edit_instructions.strip():
                await websocket.send_text(json.dumps({
                    "error": "Both current content and edit instructions are required"
                }, default=safe_json_serialize))
                continue
            
            try:
                full_response = ""
                async for chunk in stream_ai_edit(current_content, edit_instructions):
                    if chunk:
                        full_response += chunk
                        await websocket.send_text(json.dumps({
                            "chunk": chunk,
                            "connection_id": connection_id
                        }, default=safe_json_serialize))
                        await asyncio.sleep(0.1)
                
                await websocket.send_text(json.dumps({
                    "final": {
                        "edited_content": full_response,
                        "connection_id": connection_id
                    }
                }, default=safe_json_serialize))
                logger.info(f"Streaming AI edit completed for connection: {connection_id}")
            
            except Exception as e:
                logger.error(f"Error in streaming AI edit for connection {connection_id}: {e}")
                await websocket.send_text(json.dumps({
                    "error": f"Error editing content: {str(e)}",
                    "connection_id": connection_id
                }, default=safe_json_serialize))
    
    except WebSocketDisconnect:
        logger.info(f"AI Edit WebSocket disconnected: {connection_id}")
        websocket_connections.pop(connection_id, None)
    except Exception as e:
        logger.error(f"AI Edit WebSocket error: {e}")
        await websocket.send_text(json.dumps({
            "error": f"WebSocket error: {str(e)}",
            "connection_id": connection_id
        }, default=safe_json_serialize))
        websocket_connections.pop(connection_id, None)

@app.exception_handler(HTTPException)
async def http_exception_handler(request, exc):
    return JSONResponse(
        status_code=exc.status_code,
        content={
            "error": exc.detail,
            "status_code": exc.status_code,
            "timestamp": datetime.now().isoformat()
        }
    )

@app.exception_handler(Exception)
async def general_exception_handler(request, exc):
    logger.error(f"Unhandled exception: {exc}")
    return JSONResponse(
        status_code=500,
        content={
            "error": "Internal server error",
            "status_code": 500,
            "timestamp": datetime.now().isoformat()
        }
    )

@app.on_event("startup")
async def startup_event():
    logger.info("Agentic AI SRS Generator starting up...")
    initialize_groq_llms()
    logger.info(f"Groq LLMs available: {len(llms) > 0}")
    logger.info(f"OCR available: {ocr is not None}")
    logger.info(f"python-magic available: {MAGIC_AVAILABLE}")
    
    for directory in [UPLOAD_DIR, OUTPUT_DIR, TEMP_DIR]:
        directory.mkdir(exist_ok=True)
    
    logger.info("Agentic AI SRS Generator started successfully")

@app.on_event("shutdown")
async def shutdown_event():
    logger.info("Agentic AI SRS Generator shutting down...")
    try:
        for connection_id, websocket in websocket_connections.items():
            try:
                await websocket.close()
            except:
                pass
        websocket_connections.clear()
        for file_path in TEMP_DIR.glob("*"):
            if file_path.is_file():
                file_path.unlink()
    except Exception as e:
        logger.error(f"Error during shutdown: {e}")
    logger.info("Agentic AI SRS Generator shutdown complete")

documents_store: Dict[str, Dict[str, Any]] = {}
active_sessions: Dict[str, Dict[str, Any]] = {}
websocket_connections: Dict[str, WebSocket] = {}
session_workflows: Dict[str, Dict[str, Any]] = {}
session_progress: Dict[str, Dict[str, Any]] = {}
document_knowledge_base: Dict[str, Dict[str, Any]] = {}

class DocumentLearningAgent:
    @staticmethod
    async def analyze_document(document_content: str, document_metadata: Dict[str, Any]) -> Dict[str, Any]:
        llm_tuple = get_next_llm()
        if not llm_tuple:
            return {"error": "LLM not available"}
        llm, _ = llm_tuple
        
        prompt = ChatPromptTemplate.from_messages([
            ("system", """You are analyzing a document to extract CONTEXTUAL KNOWLEDGE for SRS generation support. This document will serve as REFERENCE MATERIAL only - not as source requirements.
            
            Extract these contextual insights:
            1. Business Domain: What industry/domain does this relate to?
            2. Technical Architecture Patterns: Common technical approaches used
            3. Stakeholder Types: Typical users and roles mentioned
            4. Quality Attributes: Performance, security, usability standards referenced
            5. Integration Patterns: Common integrations or interfaces
            6. Compliance Context: Regulations, standards, or governance mentioned
            7. Common Terminology: Domain-specific terms and their meanings
            8. Best Practices: Any recommended approaches or methodologies
            
            IMPORTANT: Focus on extracting CONTEXT and PATTERNS, not specific requirements. This analysis helps understand HOW to write better SRS documents, not WHAT requirements to include.
            
            Output as structured markdown with clear sections."""),
            ("human", "Document Type: {doc_type}\nDocument Name: {doc_name}\n\nDocument Content to Analyze for Context:\n{content}\n\nExtract contextual knowledge (NOT specific requirements):")
        ])
        parser = StrOutputParser()
        
        try:
            analysis = await invoke_llm_with_retry(llm, prompt, parser, {
                "doc_type": document_metadata.get("file_type", "unknown"),
                "doc_name": document_metadata.get("filename", "unknown"),
                "content": document_content[:8000]  # Limit content size
            })
            return {
                "document_id": document_metadata.get("id"),
                "filename": document_metadata.get("filename"),
                "analysis": analysis,
                "analyzed_at": datetime.now(),
                "content_preview": document_content[:1000],
                "metadata": document_metadata,
                "knowledge_type": "contextual_reference"  # Mark as context, not requirements
            }
        except Exception as e:
            logger.error(f"Error in document analysis: {e}")
            return {"error": str(e)}
    
    @staticmethod
    def store_document_knowledge(document_id: str, analysis: Dict[str, Any]) -> None:
        document_knowledge_base[document_id] = analysis
        logger.info(f"Stored knowledge for document: {document_id}")
    
    @staticmethod
    def get_relevant_knowledge(query: str, limit: int = 3) -> List[Dict[str, Any]]:
        relevant_docs = []
        query_lower = query.lower()
        
        for doc_id, knowledge in document_knowledge_base.items():
            relevance_score = 0
            analysis_text = str(knowledge.get("analysis", "")).lower()
            for word in query_lower.split():
                if word in analysis_text:
                    relevance_score += 1
            if relevance_score > 0:
                knowledge_copy = copy.deepcopy(knowledge)
                if knowledge_copy.get("analyzed_at"):
                    knowledge_copy["analyzed_at"] = knowledge_copy["analyzed_at"].isoformat()
                relevant_docs.append({
                    "document_id": doc_id,
                    "knowledge": knowledge_copy,
                    "relevance_score": relevance_score
                })
        
        relevant_docs.sort(key=lambda x: x["relevance_score"], reverse=True)
        return relevant_docs[:limit]
    
    @staticmethod
    def get_knowledge_summary() -> Dict[str, Any]:
        return {
            "total_documents": len(document_knowledge_base),
            "documents": {
                doc_id: {
                    "filename": knowledge.get("filename", "Unknown"),
                    "analyzed_at": knowledge.get("analyzed_at").isoformat() if knowledge.get("analyzed_at") else None,
                    "preview": knowledge.get("content_preview", "")[:200]
                }
                for doc_id, knowledge in document_knowledge_base.items()
            }
        }

if __name__ == "__main__":
    import uvicorn
    
    if not any(API_KEYS):
        print("WARNING: No GROQ_API_KEYS provided")
        print("Please add your Groq API keys to the API_KEYS list")
    
    print("Starting Agentic AI SRS Generator with Groq LLM and LangGraph...")
    uvicorn.run(
        socket_app,
        host="0.0.0.0",
        port=8000,
        log_level="info"
    )